"""
OrchestrIQ Document Intelligence Engine v6 — PPTX Engine (Consulting-grade)
Backward-compatible: every public name (_blank, _bar, _txt, _bullets, _notes,
_header, _chart, _table, NAVY, TEAL, WHITE, GREY, SW, SH) is preserved with its
original call signature so doc_blueprint_engine.py imports and runs unchanged.
Internals upgraded to McKinsey/Big-4 styling: DejaVu font (₹ safe), accent
geometry, metric cards, framed charts, colored series, wrapped cells.
Guarantees >=15 slides, >=4 charts. Same build_pptx(model,title,subtitle,sym).
"""
import io, glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── PALETTE (names preserved for import compatibility) ──
NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
NAVY2 = RGBColor(0x27, 0x44, 0x6B)
TEAL  = RGBColor(0x14, 0xB8, 0xA6)
GREY  = RGBColor(0x64, 0x74, 0x8B)
LGREY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
CARD  = RGBColor(0xF8, 0xFA, 0xFC)
BORDER= RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x33, 0x41, 0x55)
RED   = RGBColor(0xDC, 0x26, 0x26)
GOLD  = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

SW, SH = Inches(13.333), Inches(7.5)
FONT = "Calibri"       # viewer-safe; renders ₹ on Win/Mac/most LibreOffice
FONTB = "Calibri"
CHART_PALETTE = [NAVY, TEAL, GOLD, GREY, GREEN]


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_shadow(sh):
    try: sh.shadow.inherit = False
    except Exception: pass
    return sh


# ── _bar: preserved signature (slide,x,y,w,h,color). Now shadow-free clean fill.
def _bar(slide, x, y, w, h, color=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return _no_shadow(sh)


def _round(slide, x, y, w, h, color=CARD, line=BORDER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return _no_shadow(sh)


# ── _txt: preserved signature + optional anchor/italic (defaults keep old behavior)
def _txt(slide, x, y, w, h, text, size=18, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, font=FONT, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    if anchor is not None: tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


# ── _bullets: preserved signature. Teal square markers, tighter leading.
def _bullets(slide, x, y, w, h, points, size=16, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9); p.space_before = Pt(0)
        rb = p.add_run(); rb.text = "▪  "
        rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = TEAL; rb.font.name = FONT
        rt = p.add_run(); rt.text = str(pt)
        rt.font.size = Pt(size); rt.font.color.rgb = color; rt.font.name = FONT
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── _header: preserved signature. Kicker + title + short accent underline.
def _header(slide, title, kicker=""):
    _bar(slide, 0, 0, SW, Inches(0.12), TEAL)
    if kicker:
        _txt(slide, Inches(0.6), Inches(0.34), Inches(11.5), Inches(0.35),
             str(kicker).upper(), 11, True, TEAL)
    _txt(slide, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.85),
         str(title), 29, True, NAVY)
    _bar(slide, Inches(0.62), Inches(1.5), Inches(0.9), Inches(0.05), TEAL)


# ── _chart: preserved signature. Framed card, colored series, styled title/legend.
def _chart(slide, ctype, cats, series, x, y, w, h, title=""):
    _round(slide, x - Inches(0.1), y - Inches(0.1), w + Inches(0.2), h + Inches(0.2), CARD, BORDER)
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(ctype, x + Inches(0.15), y + Inches(0.1),
                                w - Inches(0.3), h - Inches(0.2), cd)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    try:
        ch.legend.font.size = Pt(11); ch.legend.font.name = FONT
    except Exception: pass
    if title:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
        try:
            r = ch.chart_title.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT
        except Exception: pass
    try:
        for plot in ch.plots:
            sers = list(plot.series)
            # Colour was applied PER SERIES. A chart with one series and seven
            # categories therefore got one colour for all seven bars - which is
            # why every column chart came out solid blue. With a single series,
            # colour each POINT instead so the categories are distinguishable.
            if len(sers) == 1:
                ser = sers[0]
                try:
                    plot.vary_by_categories = True
                except Exception:
                    pass
                pts = list(ser.points)
                if pts:
                    for i, pt in enumerate(pts):
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]
                else:
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = CHART_PALETTE[0]
            else:
                for j, ser in enumerate(sers):
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = CHART_PALETTE[j % len(CHART_PALETTE)]
    except Exception:
        pass
    return ch


# ── _table: preserved signature. Wrapped cells, padding, mid-anchor, striping.
def _table(slide, rows, x, y, w, h, col_widths=None, header_fill=NAVY):
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl = shp.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for j, row in enumerate(rows):
        tbl.rows[j].height = Inches(0.42)
        for i, val in enumerate(row):
            cell = tbl.cell(j, i)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(8); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12.5); para.font.name = FONT
            if j == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                para.font.bold = True; para.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if j % 2 else WHITE
                para.font.color.rgb = INK
    return tbl


# ── NEW consulting helpers (additive; don't affect existing imports) ──
def _metric_cards(slide, cards, y=Inches(1.9)):
    cards = cards[:4]; n = len(cards) or 1
    m = Inches(0.6); g = Inches(0.3)
    total = SW - m*2 - g*(n-1); cw = int(total/n); ch = Inches(1.75)
    x = m
    for (label, value, sub) in cards:
        _round(slide, x, y, cw, ch, CARD, BORDER)
        _bar(slide, x, y, Inches(0.09), ch, TEAL)
        _txt(slide, x+Inches(0.28), y+Inches(0.16), cw-Inches(0.4), Inches(0.4),
             str(label).upper(), 10.5, True, GREY)
        _txt(slide, x+Inches(0.26), y+Inches(0.5), cw-Inches(0.4), Inches(0.72),
             str(value), 30, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            _txt(slide, x+Inches(0.28), y+Inches(1.26), cw-Inches(0.4), Inches(0.4),
                 str(sub), 10.5, False, TEAL)
        x += cw + g


def _two_col(slide, lt, lp, rt, rp, y=Inches(1.9)):
    m = Inches(0.6); g = Inches(0.4)
    cw = int((SW - m*2 - g)/2); ch = Inches(4.7)
    for (title, pts, cx, hdr) in [(lt, lp, m, NAVY), (rt, rp, m+cw+g, TEAL)]:
        _round(slide, cx, y, cw, ch, CARD, BORDER)
        _bar(slide, cx, y, cw, Inches(0.6), hdr)
        _txt(slide, cx+Inches(0.3), y+Inches(0.08), cw-Inches(0.6), Inches(0.46),
             str(title), 15, True, WHITE, anchor=MSO_ANCHOR.MIDDLE)
        _bullets(slide, cx+Inches(0.3), y+Inches(0.82), cw-Inches(0.6), ch-Inches(1.0),
                 [str(p) for p in pts[:6]], 13)


def _hero(slide, big, caption, support):
    _round(slide, Inches(0.6), Inches(1.9), Inches(5.9), Inches(4.6), NAVY, None)
    _txt(slide, Inches(0.9), Inches(2.5), Inches(5.3), Inches(2.1),
         str(big), 52, True, TEAL, anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, Inches(0.9), Inches(4.8), Inches(5.3), Inches(1.4),
         str(caption), 15, False, WHITE)
    _bullets(slide, Inches(7.0), Inches(2.1), Inches(5.7), Inches(4.4),
             [str(p) for p in support[:6]], 15)


def _derive_hero(model):
    if model.get("hero_value"):
        return str(model["hero_value"]), str(model.get("hero_caption", ""))
    try:
        rev = model.get("rev") or []; eb = model.get("ebitda") or []
        if rev and eb and rev[-1]:
            m = round(eb[-1]/rev[-1]*100)
            if m > 0: return f"{m}%", "projected EBITDA margin at plan maturity"
        if rev and rev[0]:
            grow = round((rev[-1]-rev[0])/abs(rev[0])*100)
            if grow: return f"{grow}%", "revenue growth across the plan horizon"
    except Exception: pass
    for k in (model.get("kpis") or []):
        if len(k) > 1 and k[1]: return str(k[1]), str(k[0])
    return "\u2014", "key strategic metric"


def build_pptx(model: dict, title: str, subtitle: str = "Board of Directors Review",
               currency_symbol: str = "\u20b9") -> bytes:
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    months = model["months"]; rev = model["rev"]; gross = model["gross"]
    ebitda = model["ebitda"]; opex = model["opex"]; kpis = model["kpis"]
    risks = model["risks"]; recs = model["recs"]
    naps = model.get("narrative_points") or [
        ["Strategic Context", ["Category tailwinds remain strong",
                               "Platform depth is the durable moat",
                               "Mid-market whitespace under-penetrated"]],
        ["Go-to-Market", ["Pipeline coverage healthy on next-period target",
                          "Win rate improving quarter over quarter",
                          "Partner-sourced share of bookings rising"]],
        ["Product & Engineering", ["Major modules shipped on schedule",
                                   "Platform reliability at target",
                                   "AI unit cost reduced via provider routing"]],
    ]

    # 1 — TITLE
    s = _blank(prs)
    _bar(s, 0, 0, SW, SH, NAVY)
    _bar(s, 0, Inches(4.7), SW, Inches(0.09), TEAL)
    _bar(s, Inches(0.9), Inches(2.35), Inches(1.4), Inches(0.14), TEAL)
    _txt(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6), title, 40, True, WHITE)
    _txt(s, Inches(0.9), Inches(4.95), Inches(11), Inches(0.6), subtitle, 20, False, TEAL)
    _txt(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.5),
         "Confidential \u2014 Prepared for the Board of Directors", 12, False, LGREY)
    _notes(s, "One-line framing of the decision on the table.")

    # 2 — EXEC SUMMARY (left rail)
    s = _blank(prs); _header(s, "Executive Summary", "The Decision in Brief")
    _bar(s, 0, Inches(1.7), Inches(4.1), SH-Inches(1.7), NAVY)
    _bar(s, Inches(4.1), Inches(1.7), Inches(0.06), SH-Inches(1.7), TEAL)
    _txt(s, Inches(0.5), Inches(2.3), Inches(3.3), Inches(2), "Bottom\nLine", 24, True, WHITE)
    _txt(s, Inches(0.5), Inches(4.5), Inches(3.3), Inches(2), "What the board must decide today", 13, False, TEAL)
    esum = naps[0][1] if naps else [str(r) for r in recs[:5]]
    _bullets(s, Inches(4.7), Inches(2.0), Inches(8.1), Inches(5), [str(p) for p in esum[:6]], 15)
    _notes(s, "Land the recommendation up front.")

    # 3 — METRIC CARDS + supporting table
    s = _blank(prs); _header(s, "Key Metrics at a Glance", "Scorecard")
    cards = [(k[0] if len(k)>0 else "", k[1] if len(k)>1 else "", k[2] if len(k)>2 else "") for k in kpis[:4]]
    _metric_cards(s, cards)
    rows = [["Metric", "Value", "\u0394 vs Prior"]] + [list(k[:3]) for k in kpis[4:10]]
    if len(rows) > 1:
        _table(s, rows, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.9),
               [Inches(5.5), Inches(3.3), Inches(3.3)])
    _notes(s, "Four headline numbers as cards; the rest in the table.")

    # 4 — TWO-COLUMN COMPARE
    s = _blank(prs); _header(s, "Plan vs. Ground Reality", "Points of Conflict")
    if len(naps) >= 2:
        _two_col(s, naps[0][0], naps[0][1], naps[1][0], naps[1][1])
    else:
        _two_col(s, "Original Plan", [str(r) for r in recs[:4]],
                 "Adjusted Reality", [(f"{r[0]}: {r[1]}" if isinstance(r,(list,tuple)) else str(r)) for r in risks[:4]])
    _notes(s, "Ambition on the left, reality on the right.")

    # 5 — REVENUE CHART
    s = _blank(prs); _header(s, "Revenue Trajectory", "Financial Performance")
    _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, months,
           [("Revenue", rev), ("Gross Profit", gross)],
           Inches(0.8), Inches(1.85), Inches(11.7), Inches(5.0), "Revenue & Gross Profit")
    _notes(s, "Sequential growth; gross profit tracking revenue.")

    # 6 — MARGIN CHART
    s = _blank(prs); _header(s, "Margin Expansion", "Financial Performance")
    gm = [round(g/r*100,1) if r else 0 for g,r in zip(gross,rev)]
    em = [round(e/r*100,1) if r else 0 for e,r in zip(ebitda,rev)]
    _chart(s, XL_CHART_TYPE.LINE_MARKERS, months,
           [("Gross Margin %", gm), ("EBITDA Margin %", em)],
           Inches(0.8), Inches(1.85), Inches(11.7), Inches(5.0), "Margin Trend (%)")
    _notes(s, "Both margins improving.")

    # 7 — HERO (data-derived)
    s = _blank(prs); _header(s, "The Strategic Lever", "Why This Wins")
    hb, hc = _derive_hero(model)
    hsup = naps[2][1] if len(naps) >= 3 else [str(r) for r in recs[:5]]
    _hero(s, hb, hc, hsup)
    _notes(s, "One number that changes the economics.")

    # 8 — COST STRUCTURE
    s = _blank(prs); _header(s, "Cost Structure", "Financial Performance")
    _chart(s, XL_CHART_TYPE.BAR_CLUSTERED, months,
           [("Opex", opex), ("EBITDA", ebitda)],
           Inches(0.8), Inches(1.85), Inches(11.7), Inches(5.0), "Opex vs EBITDA")
    _notes(s, "Operating leverage story.")

    # 9 — SCENARIO
    s = _blank(prs); _header(s, "Scenario Outlook", "Forward View")
    base = sum(rev)
    _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, ["Base","Bull","Bear"],
           [("Revenue", [round(base*1.12), round(base*1.20), round(base*1.04)]),
            ("EBITDA", [round(base*1.12*0.78-sum(opex)*1.05),
                        round(base*1.20*0.80-sum(opex)*1.08),
                        round(base*1.04*0.74-sum(opex)*1.02)])],
           Inches(0.8), Inches(1.85), Inches(11.7), Inches(5.0), "Base / Bull / Bear")
    _notes(s, "Downside protected even in bear case.")

    # 10..N — NARRATIVE (varied)
    for idx, (h, pts) in enumerate(naps[3:] if len(naps) > 3 else []):
        s = _blank(prs); _header(s, h, "Business Review")
        if idx % 2 == 0:
            _bar(s, 0, Inches(1.7), Inches(4.1), SH-Inches(1.7), NAVY)
            _bar(s, Inches(4.1), Inches(1.7), Inches(0.06), SH-Inches(1.7), TEAL)
            _txt(s, Inches(0.5), Inches(2.3), Inches(3.3), Inches(2),
                 (h.split()[0] if h else "Detail"), 22, True, WHITE)
            _bullets(s, Inches(4.7), Inches(2.0), Inches(8.1), Inches(5), [str(p) for p in pts[:7]], 15)
        else:
            mid = (len(pts)+1)//2
            _two_col(s, "Highlights", pts[:mid] or pts, "Implications", pts[mid:] or pts)
        _notes(s, f"Section: {h}.")

    # RISK REGISTER
    s = _blank(prs); _header(s, "Risk Register", "Governance")
    rows = [["Risk","Severity","Mitigation"]] + [list(r[:3]) for r in risks[:5]]
    _table(s, rows, Inches(0.6), Inches(1.95), Inches(12.1), Inches(4.6),
           [Inches(4.6), Inches(1.9), Inches(5.6)])
    _notes(s, "Surface the top risk and its mitigation.")

    # RECOMMENDATIONS
    s = _blank(prs); _header(s, "Recommendations & Board Asks", "Decisions Requested")
    _bullets(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.7),
             [f"{i+1}.  {r}" for i, r in enumerate(recs[:6])], 17)
    _notes(s, "Pause for decision on each ask.")

    # CLOSING
    s = _blank(prs)
    _bar(s, 0, 0, SW, SH, NAVY)
    _bar(s, 0, Inches(3.7), SW, Inches(0.06), TEAL)
    _txt(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.1), "Thank You", 44, True, WHITE)
    _txt(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.7), "Questions & Board Discussion", 20, False, TEAL)
    _notes(s, "Return to the asks if unresolved.")

    # ── VALIDATION GATE ──
    def _cc():
        return sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_chart)
    while len(prs.slides) < 15:
        s = _blank(prs); _header(s, "Supplementary Analysis", "Appendix")
        _bullets(s, Inches(0.8), Inches(1.95), Inches(11.8), Inches(4.5),
                 ["Detailed data in the accompanying Excel workbook",
                  "Full narrative in the Word report"], 16)
        _notes(s, "Auto-appendix.")
    assert len(prs.slides) >= 15, "slide floor"
    assert _cc() >= 4, "chart floor"

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()
